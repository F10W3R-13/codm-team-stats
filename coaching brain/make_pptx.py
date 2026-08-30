"""
Post-Loss Speech (Point Major) — PPT generator.
Astryx neutral tone, keyword-driven, English main.
Run: python make_pptx.py
Output: 복기연설_PPT.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# --- Astryx neutral palette (from AGENTS.md base.html tokens) ---
BG          = RGBColor(0xFA, 0xFA, 0xFA)   # --bg
SURFACE     = RGBColor(0xFF, 0xFF, 0xFF)   # --card
TEXT        = RGBColor(0x26, 0x26, 0x26)   # --text
TEXT_2      = RGBColor(0x53, 0x53, 0x53)   # --text-2
MUTED       = RGBColor(0x8A, 0x8A, 0x8A)   # --muted
BORDER      = RGBColor(0xE5, 0xE5, 0xE5)   # --border
BORDER_STRONG = RGBColor(0xC8, 0xC8, 0xC8)
ACCENT      = RGBColor(0x26, 0x26, 0x26)   # --accent (black, UI)
HP          = RGBColor(0xF9, 0x73, 0x16)   # --hp (orange, functional)
SND         = RGBColor(0x8B, 0x5C, 0xF6)   # --snd (purple, functional)

# 16:9 dimensions
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Font (Pretendard preferred; falls back to Calibri/system on machines without)
FONT = "Pretendard"
FONT_FB = "Calibri"

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]  # blank layout


def _set_font(run, size, color=TEXT, bold=False, italic=False, font=FONT):
    f = run.font
    f.size = Pt(size)
    f.color.rgb = color
    f.bold = bold
    f.italic = italic
    f.name = font


def _add_text(slide, text, left, top, width, height,
              size=24, color=TEXT, bold=False, italic=False,
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    _set_font(r, size, color, bold, italic, font)
    return tb


def _rect(slide, left, top, width, height, fill, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(0.75)
    s.shadow.inherit = False
    return s


def _line(slide, left, top, width, color=BORDER_STRONG, weight=1.5):
    ln = slide.shapes.add_connector(1, left, top, left + width, top)
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    return ln


def _bg(slide, color=BG):
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, color)


def _footer(slide, idx, total, label=""):
    # bottom meta strip: small section/page marker — Astryx style
    _add_text(slide, label, Inches(0.7), Inches(7.0), Inches(8), Inches(0.3),
              size=10, color=MUTED, font=FONT)
    _add_text(slide, f"{idx:02d} / {total:02d}", Inches(11.5), Inches(7.0),
              Inches(1.5), Inches(0.3), size=10, color=MUTED, align=PP_ALIGN.RIGHT, font=FONT)


# Slide content data
TOTAL = 22

def slide_title():
    s = prs.slides.add_slide(BLANK)
    _bg(s)
    # top thin accent bar
    _rect(s, Inches(0.7), Inches(1.1), Inches(0.6), Pt(3), TEXT)
    _add_text(s, "POST-LOSS SPEECH", Inches(0.7), Inches(1.3),
              Inches(10), Inches(0.4), size=12, color=MUTED, bold=True)
    _add_text(s, "Point Major", Inches(0.7), Inches(2.6),
              Inches(12), Inches(2), size=80, color=TEXT, bold=True)
    _add_text(s, "What the loss means — and what changes today.",
              Inches(0.7), Inches(4.7), Inches(11), Inches(0.6),
              size=22, color=TEXT_2)
    _line(s, Inches(0.7), Inches(5.5), Inches(11.9))
    _add_text(s, "Team Address  ·  English", Inches(0.7), Inches(5.7),
              Inches(8), Inches(0.4), size=11, color=MUTED)
    return s


def slide_divider(num, title, subtitle=""):
    s = prs.slides.add_slide(BLANK)
    _bg(s)
    # giant section number, muted
    _add_text(s, f"{num:02d}", Inches(0.7), Inches(1.0),
              Inches(6), Inches(3.5), size=180, color=BORDER, bold=True)
    _rect(s, Inches(0.7), Inches(4.7), Inches(0.6), Pt(3), TEXT)
    _add_text(s, f"SECTION {num:02d}", Inches(0.7), Inches(4.9),
              Inches(6), Inches(0.4), size=12, color=MUTED, bold=True)
    _add_text(s, title, Inches(0.7), Inches(5.3),
              Inches(12), Inches(1.2), size=44, color=TEXT, bold=True)
    if subtitle:
        _add_text(s, subtitle, Inches(0.7), Inches(6.3),
                  Inches(12), Inches(0.6), size=18, color=TEXT_2, italic=True)
    return s


def slide_statement(text, kicker="", idx=0):
    s = prs.slides.add_slide(BLANK)
    _bg(s)
    if kicker:
        _add_text(s, kicker.upper(), Inches(0.7), Inches(1.2),
                  Inches(12), Inches(0.4), size=12, color=MUTED, bold=True)
    _add_text(s, text, Inches(0.7), Inches(2.8),
              Inches(11.9), Inches(2.5), size=54, color=TEXT, bold=True)
    return s


def slide_quote(text, attribution="", idx=0):
    s = prs.slides.add_slide(BLANK)
    _bg(s, SURFACE)
    # large accent quote mark
    _add_text(s, "\u201C", Inches(0.7), Inches(0.6),
              Inches(3), Inches(2.5), size=200, color=BORDER, bold=True)
    _add_text(s, text, Inches(1.5), Inches(2.7),
              Inches(10.5), Inches(2.5), size=38, color=TEXT, bold=True)
    if attribution:
        _line(s, Inches(1.5), Inches(5.4), Inches(0.8), BORDER_STRONG, 2)
        _add_text(s, attribution, Inches(1.5), Inches(5.6),
                  Inches(10), Inches(0.4), size=14, color=MUTED, italic=True)
    return s


def slide_bullets(items, title="", kicker="", idx=0):
    """items: list of (headline, sub) — sub may be ''"""
    s = prs.slides.add_slide(BLANK)
    _bg(s)
    if kicker:
        _add_text(s, kicker.upper(), Inches(0.7), Inches(0.8),
                  Inches(12), Inches(0.4), size=12, color=MUTED, bold=True)
    if title:
        _add_text(s, title, Inches(0.7), Inches(1.2),
                  Inches(12), Inches(0.9), size=38, color=TEXT, bold=True)
    # items as stacked rows with accent bar
    top = Inches(2.8)
    row_h = Inches(1.2)
    for i, (head, sub) in enumerate(items):
        y = top + row_h * i
        _rect(s, Inches(0.7), y + Inches(0.05), Inches(0.08), Inches(0.7), ACCENT)
        _add_text(s, head, Inches(1.0), y, Inches(11.5), Inches(0.5),
                  size=22, color=TEXT, bold=True)
        if sub:
            _add_text(s, sub, Inches(1.0), y + Inches(0.5),
                      Inches(11.5), Inches(0.5), size=15, color=TEXT_2)
    return s


def slide_closing(text, sub=""):
    s = prs.slides.add_slide(BLANK)
    _bg(s, TEXT)  # inverted: dark slide
    _rect(s, Inches(0.7), Inches(1.2), Inches(0.6), Pt(3), SURFACE)
    _add_text(s, "CLOSING", Inches(0.7), Inches(1.4),
              Inches(6), Inches(0.4), size=12, color=BORDER, bold=True)
    _add_text(s, text, Inches(0.7), Inches(2.6),
              Inches(12), Inches(2.5), size=72, color=SURFACE, bold=True)
    if sub:
        _add_text(s, sub, Inches(0.7), Inches(5.0),
                  Inches(11), Inches(0.6), size=20, color=BORDER)
    return s


# ---- build deck ----
slide_title()

# 01
slide_divider(1, "The loss wasn\u2019t an upset.",
              "It was a bill coming due.")
slide_statement("They study us. They copy us. Everything routes through us.",
                kicker="Section 01 · Why they\u2019re catching up")
slide_statement("Become impossible to study.",
                kicker="The only permanent edge")
slide_statement("Reach our level \u2192 predictable. Now what?",
                kicker="The question isn\u2019t the loss. It\u2019s what we do next.")

# 02
slide_divider(2, "The blowout trap.",
              "250-20 isn\u2019t proof of perfect play.")
slide_statement("We won by 200 because we won the gunfights \u2014 not because our play was right.",
                kicker="The lie we told ourselves")
slide_quote("Gunfights can be lost. The play around them can\u2019t.",
            attribution="How I judge a round")

# 03
slide_divider(3, "Everyone has a different definition of value.",
              "And the clash only shows up when fights start going wrong.")
slide_bullets(
    [
        ("Kings", "Lowest OBJ on the team. Value = safe angles, map control, damage positions."),
        ("Si", "Highest OBJ ratio. Value = fighting inside the hill."),
        ("Jason", "Same instinct \u2014 wants to fight inside."),
    ],
    title="The value clash, by name.",
    kicker="Section 03 · Concrete, not abstract"
)
slide_statement("\u201cI played it right. The team didn\u2019t follow.\u201d",
                kicker="The thought that breaks teams")

# 04
slide_divider(4, "There is no perfect play.",
              "Only team play.")
slide_quote("There is no perfect play. Only team play.",
            attribution="The one definition that matters")
slide_statement("Operator swipe costs us 5 seconds. Then we regroup, push, take it back.",
                kicker="How you absorb a counter")

# 05
slide_divider(5, "Two real problems.",
              "And I\u2019m not the first one naming them.")
slide_bullets(
    [
        ("The calls aren\u2019t there", "Gameplan exists. The second it gets fast, calls stop. \u201cToo busy\u201d is not an excuse at this level."),
        ("Muscle memory \u2260 knowledge", "Two separate systems. Knowing the right play \u2260 your hands doing it."),
    ],
    title="Why it isn\u2019t landing.",
    kicker="Section 05 · Both fixable. Both on us."
)
slide_statement("Override \u2192 feel the payoff \u2192 want to repeat \u2192 old habit fades.",
                kicker="The cycle I need you to run")
slide_quote("Even when you think you\u2019re right, follow me first.",
            attribution="Trust the process. I\u2019ll hand you the wheel when it\u2019s time.")

# 06
slide_divider(6, "The stakes \u2014 no spin.",
              "Point Major gone. 3.5x multiplier gone.")
slide_statement("No more \u201cI quit.\u201d  Not even as a joke.",
                kicker="Words come back in through your own ears")
slide_bullets(
    [
        ("Drop the nonchalant posture", "Too-cool-to-try doesn\u2019t win a World Championship."),
        ("Be honest. Diligent. Proud.", "Luminesis 2024 is the standard."),
        ("Feedback is a delivery skill", "If your logic is perfect but the listener shuts down \u2014 that\u2019s your loss."),
    ],
    title="Professional mindset, not cool.",
    kicker="Section 06"
)

# 07
slide_divider(7, "Was the loss worth anything?",
              "Yes \u2014 if we use it.")
slide_quote("The wake-up call came before the World Championship, not during it.",
            attribution="The best precaution we could have asked for")

# 08
slide_divider(8, "New practice structure.",
              "Starting today.")
slide_bullets(
    [
        ("Aim training \u2014 everyone, no exceptions", "XR is doing it. That\u2019s the baseline, not a flex."),
        ("7 PM EST scrim \u2014 gone", "Cold scrims are low-value. Quality > count."),
        ("Structured 3v3 warmup \u2014 the real change", "Controlled environment. One checkpoint/day. Recorded."),
        ("Targeted VOD review", "Map-specific, problem-specific. Not long anymore."),
    ],
    title="Four changes, locked in today.",
    kicker="Section 08"
)
slide_quote("3v3 \u2014 read the enemy, build counter muscle memory, feel a clean hill slam.",
            attribution="Why 3v3 beats a cold scrim")

# 09 + 10
slide_divider(9, "The system \u2014 Coaching Brain.",
              "And my part in this.")
slide_bullets(
    [
        ("The Coaching Brain", "Every map, every situation, one source of truth. VODs, transcripts, drills feed in."),
        ("We don\u2019t start from nothing", "Mapping and prep build on our philosophy \u2014 concrete plans underneath."),
        ("My part \u2014 the Argentina lesson", "I assumed we\u2019d win. That assumption is exactly us. Something has to change, and it has to be sustainable."),
        ("Predictable range", "I\u2019m not the erratic variable. If I step outside the framework, we talk first."),
    ],
    title="Institutional thinking + my commitment.",
    kicker="Sections 09 \u00b7 10"
)

slide_closing("Let\u2019s get to work.",
              "Don\u2019t question the intent. Run the cycle with me.")

# ---- add footers to body slides (skip title & dividers & closing) ----
all_slides = list(prs.slides)
# Slides we want footers on: statements, quotes, bullets (everything except 0, dividers, last)
divider_indices = {1, 4, 7, 11, 14, 18, 20, 24}  # approx; we'll skip by checking shapes
# Simpler: tag every slide that isn't title/divider/closing. We know indices:
no_footer = {0, len(all_slides) - 1}  # title and closing
# dividers are short on shapes; detect by the giant muted number text size
for i, sl in enumerate(all_slides):
    if i in no_footer:
        continue
    # detect divider: slide whose first big textbox has size >= 150
    is_divider = False
    for sh in sl.shapes:
        if sh.has_text_frame:
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    if r.font.size and r.font.size >= Pt(150):
                        is_divider = True
                        break
    if is_divider:
        continue
    _footer(sl, i + 1, TOTAL, label="Post-Loss Speech \u00b7 Point Major")

out = r"C:\Users\0616y\Downloads\Team management app\coaching brain\복기연설_PPT.pptx"
prs.save(out)
print(f"OK: {out}")
print(f"Slides: {len(prs.slides)}")
